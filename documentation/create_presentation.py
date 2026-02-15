"""
Expanded PowerPoint Presentation Generator for Fluffy Assistant
Creates a comprehensive technical dossier with code, workflows, and architecture.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    PRIMARY_COLOR = RGBColor(99, 102, 241)  # Indigo-500
    SECONDARY_COLOR = RGBColor(244, 63, 94)  # Rose-500
    ACCENT_COLOR = RGBColor(16, 185, 129)   # Emerald-500
    TEXT_COLOR = RGBColor(30, 41, 59)       # Slate-800
    GRAY_TEXT = RGBColor(71, 85, 105)       # Slate-600
    WHITE = RGBColor(255, 255, 255)
    
    # --- SLIDE 1: TITLE SLIDE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rectangle(slide, 0, 0, 13.33, 7.5, PRIMARY_COLOR)
    
    title = add_text(slide, "Fluffy Integrated Assistant System", 5, 2.5, 8, 1.5, 64, WHITE, bold=True)
    title.alignment = PP_ALIGN.CENTER
    
    subtitle = add_text(slide, "A Detailed Technical Analysis & System Documentation", 5, 4, 8, 0.8, 28, WHITE)
    subtitle.alignment = PP_ALIGN.CENTER
    
    # --- SLIDE 2: PROJECT OVERVIEW ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Project Overview", PRIMARY_COLOR)
    
    content = [
        "🐰 Fluffy is a lightweight, privacy-focused system monitor and security guardian.",
        "🛠️ Built for Windows using a high-performance multi-language stack (Rust, Python, TS).",
        "🛡️ Features signature-less behavioral threat detection (Guardian).",
        "🗣️ Integrated Voice Control and LLM Chat for an intelligent desktop experience.",
        "🔒 100% Local Processing - Zero data exfiltration by design."
    ]
    add_bullets(slide, content, 1.5, TEXT_COLOR)
    
    # --- SLIDE 3: SYSTEM ARCHITECTURE (HEXAGONAL) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Hexagonal Architecture", PRIMARY_COLOR)
    
    # Draw Architecture
    add_box(slide, 1, 1.5, 3.5, 4.5, "UI LAYER\n(Tauri + TS)", RGBColor(219, 234, 254))
    add_box(slide, 5, 1.5, 3.5, 4.5, "INTELLIGENCE LAYER\n(Python Brain)", RGBColor(224, 231, 255))
    add_box(slide, 9, 1.5, 3.3, 4.5, "CORE LAYER\n(Rust Engine)", RGBColor(254, 243, 199))
    
    add_arrow(slide, 4.5, 3.75, 5, 3.75)
    add_arrow(slide, 8.5, 3.75, 9, 3.75)
    
    add_text(slide, "HTTP/WS", 4.3, 3.3, 1.5, 0.5, 12, GRAY_TEXT)
    add_text(slide, "IPC (TCP)", 8.3, 3.3, 1.5, 0.5, 12, GRAY_TEXT)
    
    # --- SLIDE 4: FILE STRUCTURE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "File Structure Analysis", PRIMARY_COLOR)
    
    structure = [
        "/core - Rust Monitoring Engine (Performance/System Access)",
        "/brain - Python Intelligence Layer (Security/Logic/Memory)",
        "/ui/tauri - Desktop Shell and IPC Hub",
        "/ui/frontend - Modern Dashboard Assets (Vite/TypeScript)",
        "/ai - LLM Connectors and Intent Classifiers",
        "/voice - STT (Vosk) and TTS (Piper) Implementations",
        "/fluffy_data - Persistent storage (Memory/Baselines/History)"
    ]
    add_bullets(slide, structure, 1.5, TEXT_COLOR)
    
    # --- SLIDE 5: CORE SERVICE (RUST) - ROLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Core Service: The Rust Engine", PRIMARY_COLOR)
    
    roles = [
        "🏎️ Native precision: Collects system metrics every 100ms.",
        "📡 Telemetry Hub: Broadcasts JSON states over TCP port 9001.",
        "⚙️ Hardware Controller: Direct access to Volume and Brightness APIs.",
        "🛡️ Process Manager: Safe termination and tree traversal via 'sysinfo'.",
        "🌐 Network Monitor: Uses ETW for per-process packet tracking."
    ]
    add_bullets(slide, roles, 1.5, TEXT_COLOR)
    
    # --- SLIDE 6: CORE CODE: TELEMETRY LOOP ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Code Deep-Dive: Rust Telemetry", PRIMARY_COLOR)
    
    code = """
loop {
    system.refresh_cpu_all();
    system.refresh_processes(ProcessesToUpdate::All);
    
    let stats = transform_to_fluffy_json(system);
    let message = json!({ 
        "type": "telemetry", 
        "data": stats 
    }).to_string();
    
    ipc_server.broadcast(message); // Port 9001
    thread::sleep(Duration::from_millis(2000));
}
    """
    add_code_block(slide, code, 1.5, 0.5, 8, 5)
    
    explanation = [
        "• refreshes only necessary components to save CPU.",
        "• serializes to a standardized 'Fluffy' schema.",
        "• uses thread-safe broadcast to multi-client Brains."
    ]
    add_bullets(slide, explanation, 1.5, TEXT_COLOR, 9, 3.5)
    
    # --- SLIDE 7: BRAIN SERVICE (PYTHON) - ROLE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Brain Service: The Python Intelligence", PRIMARY_COLOR)
    
    roles = [
        "🧠 Semantic Monitor: Converts raw metrics into health signals.",
        "🗄️ State Manager: Thread-safe repository of current system health.",
        "🛡️ Guardian Engine: The primary behavioral analysis pipeline.",
        "🔌 Flask API: Bridge between the native UI and the internal logic.",
        "🎙️ Audio Manager: Coordinates STT input and TTS feedback."
    ]
    add_bullets(slide, roles, 1.5, TEXT_COLOR)
    
    # --- SLIDE 8: THE GUARDIAN: 4 PILLARS ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Security Guardian: Detection DNA", PRIMARY_COLOR)
    
    pillars = [
        "1. Path Integrity: Execution from %TEMP% flagged immediately (+30).",
        "2. Resource Spikes: CPU >3x baseline detected via EMA analysis (+20).",
        "3. Child Spawning: Rapid creation of CLI sub-processes tracked (+25).",
        "4. Persistence: Monitoring 'Run' registry keys and Startup folders (+40).",
        "🚀 Learning Phase: A 5-minute baseline established on first startup."
    ]
    add_bullets(slide, pillars, 1.5, TEXT_COLOR)
    
    # --- SLIDE 9: CODE: PERFORMANCE EMA ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Algorithm: Adaptive Learning (EMA)", PRIMARY_COLOR)
    
    code = """
# math.py logic
def update_baseline(old_avg, new_val, alpha=0.01):
    return (alpha * new_val) + ((1 - alpha) * old_avg)

# Guardian update
baseline['avg_cpu'] = update_baseline(
    baseline['avg_cpu'], 
    current_cpu
)
    """
    add_code_block(slide, code, 1.5, 1, 7, 4)
    
    explanation = [
        "• alpha=0.01 provides 'long-term memory'.",
        "• prevents false positives from transient spikes.",
        "• adapts to user's daily software habits."
    ]
    add_bullets(slide, explanation, 1.5, TEXT_COLOR, 8.5, 4)
    
    # --- SLIDE 10: USER INTERFACE (TAURI/TS) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "UI Layer: Modern Native Experience", PRIMARY_COLOR)
    
    features = [
        "✨ Windows Native performance via Tauri.",
        "📊 Dynamic charting (EMA-smoothed) for CPU/RAM visualization.",
        "🌳 Process Tree: Real-time parent-child hierarchy display.",
        "💬 Interactive FAB: Floating chat button for instant AI access.",
        "🎨 Glassmorphism design system with responsive layouts."
    ]
    add_bullets(slide, features, 1.5, TEXT_COLOR)
    
    # --- SLIDE 11: WORKFLOW: TELEMETRY PIPELINE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Workflow: Telemetry Data Pipeline", PRIMARY_COLOR)
    
    steps = [
        "1. Core collects metrics from OS kernel.",
        "2. Core broadcasts JSON via TCP Socket.",
        "3. Brain receives & executes Guardian logic.",
        "4. Brain updates State Store (Flask).",
        "5. UI polls State every 1s for immediate display."
    ]
    
    # Draw Flow
    y = 2.0
    for i, step in enumerate(steps):
        add_box(slide, 2.5, y, 8, 0.6, step, RGBColor(241, 245, 249))
        if i < len(steps)-1:
            add_arrow(slide, 6.5, y+0.6, 6.5, y+0.9)
        y += 0.9
        
    # --- SLIDE 12: MEMORY & PERSISTENCE ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Persistent Memory System", PRIMARY_COLOR)
    
    details = [
        "💾 JSON Long-Term Store: Saves preferences, trusted PIDs, and identities.",
        "🧠 Multi-Session Context: Remembers previous chat topics and actions.",
        "🛡️ Trusted Whitelist: Processes marked as safe persist across reboots.",
        "📝 Session Buffer: Tracks dangerous actions for delayed confirmation."
    ]
    add_bullets(slide, details, 1.5, TEXT_COLOR)
    
    # --- SLIDE 13: THE INTERRUPT SYSTEM ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Safety Control: Interrupt Commands", PRIMARY_COLOR)
    
    commands = [
        'Keywords: "Stop", "Cancel", "Abort", "Shut up", "Nevermind".',
        "⚡ Immediate Cancellation: Halts TTS playback in <50ms.",
        "🛑 Logic Reset: Clears all pending intents and confirmation prompts.",
        "🔔 Visual Feedback: Toast notifications confirm the interrupt."
    ]
    add_bullets(slide, commands, 1.5, TEXT_COLOR)
    
    # --- SLIDE 14: AI & VOICE INTEGRATION ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Natural Language Interface", PRIMARY_COLOR)
    
    stack = [
        "🎤 STT: Vosk Offline Engine (privacy-centric).",
        "🔊 TTS: Piper Neural Voices (highly realistic).",
        "🤖 LLM Support: OpenAI, Claude, Groq, and Ollama (Local).",
        "🧪 Intent Mapping: Maps speech to 15+ internal system commands."
    ]
    add_bullets(slide, stack, 1.5, TEXT_COLOR)
    
    # --- SLIDE 15: SELF-IMPROVEMENT LOOP ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The Self-Improvement Cycle", PRIMARY_COLOR)
    
    logic = [
        "👁️ Observer: Tracks user failures or 'I don't know' responses.",
        "🏗️ Architect: Suggests new Python extensions to handle new commands.",
        "🔨 Generator: Auto-writes and installs extensions into /brain/extensions.",
        "📦 Dynamic Loading: New features become active without a system restart."
    ]
    add_bullets(slide, logic, 1.5, TEXT_COLOR)
    
    # --- SLIDE 16: SYSTEM NORMALIZATION ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The 'Normalize' Feature", PRIMARY_COLOR)
    
    actions = [
        "🔊 Sound: Resets volume to 50% for optimal environment.",
        "☀️ Visuals: Forces brightness to 70-80% to avoid eye strain.",
        "🧹 Disk: Purges Windows /Temp folders of bloated logs.",
        "🔍 Security: Triggers a deep Guardian scan for all active PIDs.",
        "🎯 Goal: Bring messy systems back to a pristine baseline."
    ]
    add_bullets(slide, actions, 1.5, TEXT_COLOR)
    
    # --- SLIDE 17: CONCLUSION ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Conclusion & Future", PRIMARY_COLOR)
    
    summary = [
        "✅ Technical Excellence: Multi-service architecture for speed and safety.",
        "✅ Privacy First: No dependencies on third-party cloud data mining.",
        "✅ Future: Cross-platform support and decentralized P2P telemetry.",
        "🐰 Fluffy: Where System Monitoring meets Personal Intelligence."
    ]
    add_bullets(slide, summary, 1.5, TEXT_COLOR)
    
    # --- SLIDE 18: THANK YOU ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rectangle(slide, 0, 0, 13.33, 7.5, PRIMARY_COLOR)
    add_text(slide, "Thank You!", 5, 3, 8, 1, 72, WHITE, bold=True).alignment = PP_ALIGN.CENTER
    add_text(slide, "Questions? | documentation/agent.md", 5, 4.5, 8, 0.5, 24, WHITE).alignment = PP_ALIGN.CENTER

    # Save
    path = r'C:/Users/sudip/OneDrive/Desktop/webProjects/FluffyAssistent/documentation/presentation.pptx'
    prs.save(path)
    print(f"✅ Detailed technical presentation created at {path}")

# --- HELPERS ---

def add_title(slide, text, color):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.8))
    frame = textbox.text_frame
    frame.text = text
    p = frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = color

def add_bullets(slide, items, y_start, color, x=1, width=11):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y_start), Inches(width), Inches(5))
    frame = textbox.text_frame
    frame.word_wrap = True
    for i, item in enumerate(items):
        p = frame.add_paragraph() if i > 0 else frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = color
        p.space_after = Pt(12)

def add_code_block(slide, code, y, x, w, h):
    # Dark background for code
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(15, 23, 42) # Slate-900
    rect.line.fill.background()
    
    textbox = slide.shapes.add_textbox(Inches(x+0.2), Inches(y+0.2), Inches(w-0.4), Inches(h-0.4))
    frame = textbox.text_frame
    frame.text = code.strip()
    p = frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(148, 163, 184) # Slate-400

def add_box(slide, x, y, w, h, text, color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = RGBColor(148, 163, 184)
    frame = box.text_frame
    frame.text = text
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def add_arrow(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = RGBColor(148, 163, 184)
    connector.line.width = Pt(2)

def add_rectangle(slide, x, y, w, h, color):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()

def add_text(slide, text, x, y, w, h, size, color, bold=False):
    textbox = slide.shapes.add_textbox(Inches(x - w/2), Inches(y), Inches(w), Inches(h))
    frame = textbox.text_frame
    frame.text = text
    p = frame.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return p

if __name__ == "__main__":
    create_presentation()
